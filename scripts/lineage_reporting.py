# -*- coding: utf-8 -*-
"""
Created on Tue Jul 11 17:37:47 2023

@author: johan
"""

import networkx as nx                                    
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx import Presentation

def get_extreme_nodes(sg, exclude_nodes):
    """Returns the extreme nodes of a subgraph, excluding the nodes in 'exclude_nodes'"""
    extreme_nodes = [node for node in sg.nodes if sg.degree(node) == 1 and node not in exclude_nodes]
    return extreme_nodes

def create_pptx_report(subgraphs, filename):
    # Créer une présentation
    prs = Presentation()

    # Nombre maximum de lignes par slide
    max_rows_per_slide = 6

    # Calculer le nombre de colonnes nécessaires
    max_nodes = max(len(sg.nodes) for sg in subgraphs)
    max_cols = max_nodes + 2  # +2 for the start and end columns, and the node count column

    # Headers
    headers = ["Extrémité"] + [""]*(max_cols - 3) + ["Extrémité", "Nombre de variables"]

    # Boucle pour créer plusieurs slides si nécessaire
    for i in range(0, len(subgraphs), max_rows_per_slide):
        slide_subgraphs = subgraphs[i:i+max_rows_per_slide]

        # Ajouter une slide
        slide_layout = prs.slide_layouts[5]  # choisir un layout vide
        slide = prs.slides.add_slide(slide_layout)

        # Ajouter un titre
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Rapport de Data Lineage"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Créer une table
        rows, cols = len(slide_subgraphs) + 1, max_cols  # +1 for the header row
        top = Inches(1.5)  # la position de départ du haut
        left = Inches(0.2)  # la position de départ de la gauche
        width = Inches(12)  # la largeur du tableau
        height = Inches(5)  # la hauteur du tableau

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for column in table.columns:
            column.width = Inches(1.2)  # set column width

        # Remplir la table avec une taille de police plus petite
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = ''
                if r == 0:  # header row
                    cell.text = headers[c]
                elif r > 0:
                    sg = slide_subgraphs[r-1]
                    nodes = list(sg.nodes)
                    extreme_nodes = get_extreme_nodes(sg, [])  # Get all extreme nodes of the subgraph
                    if c == 0:  # first column
                        node = nodes[0]
                        file, sheet, variable = node.split("\n")
                        cell.text = f'{file}\n{sheet}\n{variable}'
                    elif c == cols - 2:  # second last column (end)
                        node = nodes[-1]
                        file, sheet, variable = node.split("\n")
                        cell.text = f'{file}\n{sheet}\n{variable}'
                    elif c == cols - 1:  # last column (count)
                        cell.text = str(len(nodes))
                    elif c > 0 and c < cols - 2:  # middle columns
                        if c-1 < len(nodes) - 2:  # there are still nodes to place in the table
                            node = nodes[c]
                            file, sheet, variable = node.split("\n")
                            cell.text = f'{file}\n{sheet}\n{variable}'
        
                # mise en forme du texte
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)  # taille de la police
                        node_as_str = '\n'.join(cell.text.split('\n')[:3])
                        if r == 0 or node_as_str in extreme_nodes or c in [0, cols-2]:  # if node is an extreme node or in first and second last column
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black color



    # Enregistrer la présentation
    prs.save(filename)

# import du graphe
G = nx.read_gexf("lineage_graph.gexf")

# Convertir le graphe dirigé en graphe non dirigé pour obtenir les sous-graphes
UG = G.to_undirected()

# Obtenir les sous-graphes
subgraphs = [G.subgraph(c) for c in nx.connected_components(UG)]

# Créer le rapport
create_pptx_report(subgraphs, "data_lineage_report.pptx")
